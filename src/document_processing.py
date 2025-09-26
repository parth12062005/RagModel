"""
Document processing functions for various file types
"""
import requests
import socket
import ssl
from urllib.parse import urlparse
from typing import List, Dict, Any, Tuple
import time

from src.config import LARGE_PDF_THRESHOLD, MAX_FILE_SIZE, UNKNOWN_FILE_TIMEOUT, UNKNOWN_FILE_MAX_SIZE, ChunkingConfig

def fast_text_only_extract(pdf_bytes: bytes) -> List[str]:
    """Ultra-fast text-only extraction for large PDFs"""
    import fitz
    
    chunks = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    for page_num, page in enumerate(doc):
        page_text = page.get_text("text").strip()
        
        if page_text and len(page_text) > 20:
            chunks.append(f"[Page {page_num + 1}]\n{page_text}")
    
    doc.close()
    return chunks

def fast_quality_pdf_extract(pdf_bytes: bytes) -> List[str]:
    """Optimized PDF extraction with comprehensive table extraction"""
    import fitz
    
    chunks = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    for page_num, page in enumerate(doc):
        # Get blocks and sort by position for better structure
        blocks = page.get_text("blocks")
        sorted_blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
        
        page_content = []
        for block in sorted_blocks:
            text = block[4].strip()
            if text and len(text) > 10:
                # Simple header detection
                if len(text) < 80 and text.isupper():
                    text = f"HEADER: {text}"
                page_content.append(text)
        
        if page_content:
            page_text = "\n\n".join(page_content)
            chunks.append(f"[Page {page_num + 1}]\n{page_text}")
        
        # Table extraction
        try:
            tables = page.find_tables()
            for table_idx, table in enumerate(tables):
                try:
                    table_data = table.extract()
                    if table_data:
                        table_text = format_table_as_text(
                            table_data, 
                            f"Page {page_num + 1}, Table {table_idx + 1}"
                        )
                        if table_text.strip():
                            chunks.append(table_text)
                            
                except Exception as e:
                    print(f"⚠️ Table extraction failed on page {page_num + 1}, table {table_idx + 1}: {e}")
        except Exception as e:
            print(f"⚠️ Table detection failed on page {page_num + 1}: {e}")
    
    doc.close()
    return chunks

def format_table_as_text(table_data: list, table_title: str) -> str:
    """Helper function to format table data as readable text"""
    if not table_data:
        return ""
    
    formatted_lines = [f"\n[{table_title}]"]
    
    for row in table_data:
        if row and any(cell for cell in row if cell):
            row_text = " | ".join(str(cell or "").strip() for cell in row)
            if row_text.strip():
                formatted_lines.append(row_text)
    
    return "\n".join(formatted_lines)

def smart_chunk_configurable(text: str, chunking_config: ChunkingConfig) -> List[str]:
    """Configurable smart chunking"""
    chunks = []
    start = 0
    chunk_size = chunking_config.chunk_size
    overlap = chunking_config.overlap
    
    while start < len(text):
        end = start + chunk_size
        
        if end >= len(text):
            chunk = text[start:].strip()
            if chunk:
                chunks.append(chunk)
            break
        
        # Find good break points (in order of preference)
        break_points = [
            text.rfind('\n\n', start + chunk_size//2, end),  # Paragraph
            text.rfind('. ', start + chunk_size//2, end),    # Sentence
            text.rfind('\n', start + chunk_size//2, end),    # Line
            text.rfind(' ', start + chunk_size//2, end)      # Word
        ]
        
        break_point = next((bp for bp in break_points if bp > -1), -1)
        
        if break_point > -1:
            chunk = text[start:break_point + 1].strip()
            next_start = break_point + 1 - overlap
        else:
            chunk = text[start:end].strip()
            next_start = end - overlap
        
        if chunk:
            chunks.append(chunk)
        
        start = max(next_start, start + 1)
    
    return chunks

def detect_file_type(url: str, content_type: str = None, content_bytes: bytes = None) -> str:
    """
    Detect file type from URL, content-type header, and file content
    """
    url_lower = url.lower()
    content_type_lower = (content_type or "").lower()
    
    # Check URL extension first
    if url_lower.endswith('.pdf'):
        return 'PDF'
    elif url_lower.endswith('.docx'):
        return 'DOCX'
    elif url_lower.endswith(('.ppt', '.pptx')):
        return 'POWERPOINT'
    elif url_lower.endswith(('.xls', '.xlsx')):
        return 'EXCEL'
    elif url_lower.endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp')):
        return 'IMAGE'
    elif url_lower.endswith(('.txt', '.text')):
        return 'TEXT'
    elif url_lower.endswith('.csv'):
        return 'CSV'
    
    # Check content-type header
    if 'pdf' in content_type_lower:
        return 'PDF'
    elif 'officedocument.wordprocessingml' in content_type_lower:
        return 'DOCX'
    elif 'officedocument.presentationml' in content_type_lower or 'powerpoint' in content_type_lower:
        return 'POWERPOINT'
    elif 'officedocument.spreadsheetml' in content_type_lower or 'excel' in content_type_lower:
        return 'EXCEL'
    elif content_type_lower.startswith('image/'):
        return 'IMAGE'
    elif 'text/plain' in content_type_lower:
        return 'TEXT'
    elif 'text/csv' in content_type_lower:
        return 'CSV'
    
    # Check file content magic bytes if available
    if content_bytes:
        if content_bytes.startswith(b'%PDF'):
            return 'PDF'
        elif content_bytes.startswith(b'PK\x03\x04'):
            # ZIP-based formats (Office documents)
            if b'word/' in content_bytes[:1024]:
                return 'DOCX'
            elif b'ppt/' in content_bytes[:1024]:
                return 'POWERPOINT'
            elif b'xl/' in content_bytes[:1024]:
                return 'EXCEL'
        elif content_bytes.startswith((b'\xff\xd8\xff', b'\x89PNG', b'GIF8')):
            return 'IMAGE'
    
    # Default to unknown
    return 'UNKNOWN'

def safe_download_with_limits(url: str, max_size_mb: int = 500, timeout_seconds: int = 30) -> Tuple[bytes, str]:
    """
    Safely download file with size and timeout limits
    Returns (content_bytes, content_type)
    """
    # First, check file size with HEAD request
    try:
        head_response = requests.head(url, timeout=10)
        content_length = head_response.headers.get('content-length')
        content_type = head_response.headers.get('content-type', '')
        
        if content_length:
            file_size = int(content_length)
            max_bytes = max_size_mb * 1024 * 1024
            
            if file_size > max_bytes:
                raise ValueError(f"File too large: {file_size/1024/1024:.1f}MB exceeds {max_size_mb}MB limit")
            
            print(f"📊 File size check: {file_size/1024/1024:.1f}MB (within {max_size_mb}MB limit)")
        
    except requests.exceptions.RequestException as e:
        print(f"⚠️ Could not check file size via HEAD request: {e}")
        content_type = ""
    
    # Download with streaming and size checking
    print(f"⬇️ Downloading with {timeout_seconds}s timeout...")
    response = requests.get(url, stream=True, timeout=timeout_seconds)
    response.raise_for_status()
    
    content_type = response.headers.get('content-type', content_type)
    content = b""
    max_bytes = max_size_mb * 1024 * 1024
    
    for chunk in response.iter_content(chunk_size=8192):
        content += chunk
        if len(content) > max_bytes:
            raise ValueError(f"Download exceeded {max_size_mb}MB limit during streaming")
    
    print(f"✅ Successfully downloaded {len(content)/1024/1024:.1f}MB")
    return content, content_type

def extract_document_content(url: str, file_type: str, content_bytes: bytes) -> List[str]:
    """
    Extract content from various document types
    """
    if file_type == 'PDF':
        # Choose extraction method based on file size
        if len(content_bytes) > LARGE_PDF_THRESHOLD:
            print(f"📄 Large PDF detected ({len(content_bytes)/1024/1024:.1f}MB), using fast text-only extraction")
            return fast_text_only_extract(content_bytes)
        else:
            print(f"📄 PDF detected ({len(content_bytes)/1024/1024:.1f}MB), using full extraction with tables")
            return fast_quality_pdf_extract(content_bytes)
    
    elif file_type == 'DOCX':
        return extract_docx_content(content_bytes)
    
    elif file_type == 'POWERPOINT':
        return extract_powerpoint_content(content_bytes)
    
    elif file_type == 'EXCEL':
        return extract_excel_content(content_bytes)
    
    elif file_type == 'IMAGE':
        return extract_image_content(content_bytes)
    
    elif file_type == 'TEXT':
        return extract_text_content(content_bytes)
    
    elif file_type == 'CSV':
        return extract_csv_content(content_bytes)
    
    else:
        return extract_unknown_content(url, content_bytes)

def extract_docx_content(content_bytes: bytes) -> List[str]:
    """Extract content from DOCX files"""
    try:
        from docx import Document
        import io
        
        doc = Document(io.BytesIO(content_bytes))
        chunks = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                chunks.append(para.text.strip())
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(row_data):
                    table_data.append(row_data)
            
            if table_data:
                table_text = format_table_as_text(table_data, "DOCX Table")
                chunks.append(table_text)
        
        return chunks
    except Exception as e:
        print(f"⚠️ DOCX extraction failed: {e}")
        return [f"Error extracting DOCX content: {str(e)}"]

def extract_powerpoint_content(content_bytes: bytes) -> List[str]:
    """Extract content from PowerPoint files"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(content_bytes))
        chunks = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                chunks.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text))
        
        return chunks
    except Exception as e:
        print(f"⚠️ PowerPoint extraction failed: {e}")
        return [f"Error extracting PowerPoint content: {str(e)}"]

def extract_excel_content(content_bytes: bytes) -> List[str]:
    """Extract content from Excel files"""
    try:
        import pandas as pd
        import io
        
        # Read all sheets
        excel_file = pd.ExcelFile(io.BytesIO(content_bytes))
        chunks = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            # Add sheet summary
            chunks.append(f"[Sheet: {sheet_name}]")
            chunks.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            
            # Add column information
            for col in df.columns:
                if df[col].dtype in ['int64', 'float64']:
                    chunks.append(f"Column '{col}': {df[col].dtype}, Mean: {df[col].mean():.2f}")
                else:
                    unique_count = df[col].nunique()
                    chunks.append(f"Column '{col}': {df[col].dtype}, Unique values: {unique_count}")
            
            # Add sample data
            sample_data = df.head(10).to_string()
            chunks.append(f"Sample data:\n{sample_data}")
        
        return chunks
    except Exception as e:
        print(f"⚠️ Excel extraction failed: {e}")
        return [f"Error extracting Excel content: {str(e)}"]

def extract_image_content(content_bytes: bytes) -> List[str]:
    """Extract text from images using OCR"""
    try:
        import pytesseract
        from PIL import Image
        import io
        
        # Load image
        image = Image.open(io.BytesIO(content_bytes))
        
        # Extract text using OCR
        text = pytesseract.image_to_string(image)
        
        if text.strip():
            return [f"[OCR Text]\n{text.strip()}"]
        else:
            return ["[Image] No text detected in image"]
    except Exception as e:
        print(f"⚠️ Image OCR failed: {e}")
        return [f"Error extracting text from image: {str(e)}"]

def extract_text_content(content_bytes: bytes) -> List[str]:
    """Extract content from text files"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                text = content_bytes.decode(encoding)
                # Split into paragraphs
                paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                return paragraphs
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, use utf-8 with errors='replace'
        text = content_bytes.decode('utf-8', errors='replace')
        return [text]
    except Exception as e:
        print(f"⚠️ Text extraction failed: {e}")
        return [f"Error extracting text content: {str(e)}"]

def extract_csv_content(content_bytes: bytes) -> List[str]:
    """Extract content from CSV files"""
    try:
        import pandas as pd
        import io
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                df = pd.read_csv(io.BytesIO(content_bytes), encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            # If all encodings fail, use utf-8 with errors='replace'
            df = pd.read_csv(io.BytesIO(content_bytes), encoding='utf-8', errors='replace')
        
        chunks = []
        chunks.append(f"[CSV Data] Rows: {len(df)}, Columns: {len(df.columns)}")
        
        # Add column information
        for col in df.columns:
            if df[col].dtype in ['int64', 'float64']:
                chunks.append(f"Column '{col}': {df[col].dtype}, Mean: {df[col].mean():.2f}")
            else:
                unique_count = df[col].nunique()
                chunks.append(f"Column '{col}': {df[col].dtype}, Unique values: {unique_count}")
        
        # Add sample data
        sample_data = df.head(10).to_string()
        chunks.append(f"Sample data:\n{sample_data}")
        
        return chunks
    except Exception as e:
        print(f"⚠️ CSV extraction failed: {e}")
        return [f"Error extracting CSV content: {str(e)}"]

def extract_unknown_content(url: str, content_bytes: bytes) -> List[str]:
    """Extract metadata from unknown file types"""
    try:
        from urllib.parse import urlparse
        
        parsed_url = urlparse(url)
        filename = parsed_url.path.split('/')[-1] if parsed_url.path else "unknown_file"
        
        chunks = []
        chunks.append(f"[Unknown File Type]")
        chunks.append(f"Filename: {filename}")
        chunks.append(f"File size: {len(content_bytes)} bytes")
        chunks.append(f"URL: {url}")
        
        # Try to extract any readable text
        try:
            text = content_bytes.decode('utf-8', errors='replace')
            if len(text.strip()) > 10:
                chunks.append(f"Extracted text (first 500 chars):\n{text[:500]}...")
        except:
            pass
        
        return chunks
    except Exception as e:
        print(f"⚠️ Unknown content extraction failed: {e}")
        return [f"Error extracting unknown content: {str(e)}"]
